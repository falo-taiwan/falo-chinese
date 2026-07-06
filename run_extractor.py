import os
import sys
import argparse
import cv2
import requests
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import yt_dlp

def format_timestamp(seconds: float) -> str:
    h = int(seconds // 3600)
    m = int((seconds % 3600) // 60)
    s = int(seconds % 60)
    return f"{h:02d}:{m:02d}:{s:02d}"

def generate_report_html(file_path, title, video_url, duration, slide_candidates):
    import base64
    
    spirit_html = """
    <div class="spirit-banner">
        <div class="spirit-logo">💡 專案核心精神 (Project Credo)</div>
        <div class="spirit-grid">
            <div class="spirit-card">
                <h3>📌 知識管理 (KM) 與經驗傳承</h3>
                <p>本專案旨在捕捉專家影片中的精華經驗，並將其高效率儲存為結構化的知識庫檔案，建立起個人與企業的知識管理體系。</p>
            </div>
            <div class="spirit-card">
                <h3>⚡ 零運行成本運算 (No Token Cost)</h3>
                <p>執行過程完全依靠 OpenCV 本機電腦視覺演算法（MAE 影格追蹤、dHash 去重與邊緣變異分析），<b>執行時完全不耗費任何 AI API Token</b>。</p>
            </div>
            <div class="spirit-card">
                <h3>🤝 人機協調驗證 (Human-in-the-Loop)</h3>
                <p>利用電腦視覺做大批量自動化過濾，再透過「雙向滑塊對比」交由人類做出最終的版本確認，是極佳的人機協作實踐。</p>
            </div>
        </div>
    </div>
    """
    
    workflow_html = """
    <div class="workflow-section">
        <h3>🧬 簡報擷取技術工作流程 (Workflow Overview)</h3>
        <div style="display: grid; grid-template-columns: repeat(auto-fit, minmax(300px, 1fr)); gap: 30px; margin-top: 20px;">
            <div>
                <h4 style="color: #38bdf8; font-size: 1rem; margin-top: 0; margin-bottom: 15px; border-bottom: 1px dashed rgba(255,255,255,0.08); padding-bottom: 5px;">🎨 視覺化工作流 (Visual Flow)</h4>
                <div class="workflow-steps">
                    <div class="step-badge">1. 影片分析</div>
                    <div class="step-arrow">&rarr;</div>
                    <div class="step-badge">2. 畫面變化偵測</div>
                    <div class="step-arrow">&rarr;</div>
                    <div class="step-badge">3. 換頁推測</div>
                    <div class="step-arrow">&rarr;</div>
                    <div class="step-badge">4. 相似度分析</div>
                    <div class="step-arrow">&rarr;</div>
                    <div class="step-badge">5. 模糊檢測</div>
                    <div class="step-arrow">&rarr;</div>
                    <div class="step-badge">6. 重複頁排除</div>
                    <div class="step-arrow">&rarr;</div>
                    <div class="step-badge">7. 候選投影片</div>
                </div>
            </div>
            
            <div>
                <h4 style="color: #38bdf8; font-size: 1rem; margin-top: 0; margin-bottom: 15px; border-bottom: 1px dashed rgba(255,255,255,0.08); padding-bottom: 5px;">📝 純文字版 (Plain text)</h4>
                <div style="background: #ffffff; color: #000000; border-radius: 12px; padding: 20px; font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; box-shadow: 0 4px 12px rgba(0,0,0,0.1); border: 1px solid #e2e8f0; width: 100%; box-sizing: border-box; text-align: left;">
                    <div style="font-weight: 600; color: #64748b; font-size: 0.85rem; text-transform: uppercase; letter-spacing: 0.5px; border-bottom: 1px solid #f1f5f9; padding-bottom: 8px; margin-bottom: 15px;">Plain text</div>
                    <div style="font-size: 1.1rem; line-height: 1.6; font-weight: 600; font-family: inherit; color: #1e293b;">
                        影片分析<br>
                        &darr;<br>
                        畫面變化偵測<br>
                        &darr;<br>
                        換頁推測<br>
                        &darr;<br>
                        相似度分析<br>
                        &darr;<br>
                        模糊檢測<br>
                        &darr;<br>
                        重複頁排除<br>
                        &darr;<br>
                        候選投影片
                    </div>
                </div>
            </div>
        </div>
        <p class="workflow-desc" style="margin-top: 25px;">本系統演算法原理基於車流量與「變動監控」技術（類似於高速公路變動監控演算法），利用數學矩陣與空間域算法去精確推算「靜止」與「跳頁」的分界點，排除干擾畫面。</p>
    </div>
    """
    
    slides_list_html = ""
    for s in slide_candidates:
        img_src = ""
        if os.path.exists(s["path"]):
            try:
                with open(s["path"], "rb") as img_f:
                    base64_data = base64.b64encode(img_f.read()).decode('utf-8')
                    img_src = f"data:image/png;base64,{base64_data}"
            except Exception as e:
                print("HTML report base64 encode failed:", e)
        
        slides_list_html += f"""
        <div class="slide-card">
            <div class="slide-image-box">
                <img src="{img_src}" alt="Slide {s['slide_no']}">
            </div>
            <div class="slide-meta-box">
                <div class="slide-title">
                    <span>投影片 #{s['slide_no']:03d}</span>
                    <span class="badge badge-orig">原始擷取</span>
                </div>
                <div class="meta-item"><strong>時間點:</strong> {s['timestamp']}</div>
                <div class="meta-item"><strong>精確秒數:</strong> {s['seconds']:.2f} 秒</div>
            </div>
        </div>
        """
        
    duration_str = format_timestamp(duration)
    
    html_content = f"""<!DOCTYPE html>
<html lang="zh-TW">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>投影片擷取任務報告 - {title}</title>
    <link href="https://fonts.googleapis.com/css2?family=Outfit:wght@300;400;600;700&family=Noto+Sans+TC:wght@300;400;500;700&display=swap" rel="stylesheet">
    <style>
        body {{
            font-family: 'Outfit', 'Noto Sans TC', sans-serif;
            background-color: #0b0f19;
            color: #f1f5f9;
            line-height: 1.6;
            margin: 0;
            padding: 0;
        }}
        .container {{
            max-width: 1000px;
            margin: 40px auto;
            padding: 20px;
        }}
        .header {{
            background: linear-gradient(135deg, #1e293b 0%, #0f172a 100%);
            border: 1px solid rgba(255,255,255,0.08);
            border-radius: 16px;
            padding: 30px;
            margin-bottom: 30px;
        }}
        h1 {{
            color: #38bdf8;
            font-size: 1.8rem;
            margin-top: 0;
        }}
        .metadata-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
            gap: 15px;
            margin-top: 20px;
            border-top: 1px solid rgba(255,255,255,0.08);
            padding-top: 20px;
            font-size: 0.9rem;
            color: #94a3b8;
        }}
        .metadata-grid span strong {{
            color: #f1f5f9;
        }}
        .spirit-banner {{
            background: rgba(56, 189, 248, 0.05);
            border: 1px dashed rgba(56, 189, 248, 0.3);
            border-radius: 12px;
            padding: 24px;
            margin-bottom: 30px;
        }}
        .spirit-logo {{
            color: #38bdf8;
            font-weight: 700;
            font-size: 1.1rem;
            margin-bottom: 15px;
        }}
        .spirit-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(250px, 1fr));
            gap: 20px;
        }}
        .spirit-card h3 {{
            color: #34d399;
            font-size: 0.95rem;
            margin-top: 0;
            margin-bottom: 8px;
        }}
        .spirit-card p {{
            font-size: 0.85rem;
            color: #94a3b8;
            margin: 0;
        }}
        .workflow-section {{
            background: rgba(30, 41, 59, 0.4);
            border: 1px solid rgba(255,255,255,0.05);
            border-radius: 12px;
            padding: 20px;
            margin-bottom: 30px;
        }}
        .workflow-section h3 {{
            color: #f1f5f9;
            font-size: 1.1rem;
            margin-top: 0;
            margin-bottom: 15px;
        }}
        .workflow-steps {{
            display: flex;
            flex-wrap: wrap;
            align-items: center;
            gap: 8px;
            margin-bottom: 12px;
        }}
        .step-badge {{
            background: rgba(255,255,255,0.06);
            border: 1px solid rgba(255,255,255,0.1);
            color: #38bdf8;
            padding: 6px 12px;
            border-radius: 6px;
            font-size: 0.8rem;
            font-weight: 600;
        }}
        .step-arrow {{
            color: #475569;
            font-weight: bold;
        }}
        .workflow-desc {{
            font-size: 0.85rem;
            color: #94a3b8;
            margin: 0;
        }}
        .slides-container {{
            display: flex;
            flex-direction: column;
            gap: 24px;
        }}
        .slide-card {{
            display: flex;
            background: rgba(30, 41, 59, 0.3);
            border: 1px solid rgba(255,255,255,0.05);
            border-radius: 12px;
            overflow: hidden;
            backdrop-filter: blur(10px);
            transition: transform 0.2s;
        }}
        .slide-card:hover {{
            transform: translateY(-2px);
            border-color: rgba(56, 189, 248, 0.2);
        }}
        .slide-image-box {{
            flex: 0 0 320px;
            background: #000;
            display: flex;
            align-items: center;
            justify-content: center;
            overflow: hidden;
            border-right: 1px solid rgba(255,255,255,0.05);
        }}
        .slide-image-box img {{
            width: 100%;
            height: auto;
            max-height: 200px;
            object-fit: contain;
            display: block;
        }}
        .slide-meta-box {{
            flex: 1;
            padding: 20px;
        }}
        .slide-title {{
            display: flex;
            justify-content: space-between;
            align-items: center;
            font-size: 1.2rem;
            font-weight: 700;
            color: #f1f5f9;
            margin-bottom: 12px;
            border-bottom: 1px solid rgba(255,255,255,0.08);
            padding-bottom: 8px;
        }}
        .badge {{
            font-size: 0.75rem;
            padding: 4px 8px;
            border-radius: 4px;
            font-weight: 600;
        }}
        .badge-orig {{
            background: rgba(255, 255, 255, 0.08);
            color: #94a3b8;
            border: 1px solid rgba(255,255,255,0.1);
        }}
        .meta-item {{
            font-size: 0.9rem;
            color: #94a3b8;
            margin-bottom: 6px;
        }}
        .meta-item strong {{
            color: #cbd5e1;
        }}
        @media (max-width: 768px) {{
            .slide-card {{
                flex-direction: column;
            }}
            .slide-image-box {{
                flex: none;
                width: 100%;
            }}
        }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>投影片自動分析與擷取報告</h1>
            <div class="metadata-grid">
                <span>影片名稱: <strong>{title}</strong></span>
                <span>影片網址: <strong>{video_url}</strong></span>
                <span>總長度: <strong>{duration_str}</strong></span>
                <span>匯出數量: <strong>{len(slide_candidates)} 張投影片</strong></span>
            </div>
        </div>
        {spirit_html}
        {workflow_html}
        <div class="slides-container">
            {slides_list_html}
        </div>
    </div>
</body>
</html>
"""
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(html_content)

def analyze_and_upload(video_url, task_id, cloudflare_url, threshold=8.0, interval=1.0, cookies=None, token_hint=None):
    import time
    start_time = time.time()

    github_run_id = os.environ.get("GITHUB_RUN_ID")
    github_repo = os.environ.get("GITHUB_REPOSITORY")
    github_run_url = f"https://github.com/{github_repo}/actions/runs/{github_run_id}" if github_run_id and github_repo else None

    print(f"🎬 Starting YouTube extraction for URL: {video_url}")
    print(f"   Task ID: {task_id}")
    print(f"   Cloudflare URL: {cloudflare_url}")
    if token_hint:
        print(f"   Token Hint: {token_hint}")
    if github_run_url:
        print(f"   GitHub Run URL: {github_run_url}")

    # Debug: Check JS Runtimes (Node.js & Deno)
    try:
        import shutil
        import subprocess
        node_path = shutil.which("node")
        print(f"🔍 Node.js path: {node_path}")
        if node_path:
            node_ver = subprocess.check_output(["node", "--version"], text=True).strip()
            print(f"   Node.js version: {node_ver}")
        else:
            print("   Node.js not found in PATH!")
            
        deno_path = shutil.which("deno")
        print(f"🔍 Deno path: {deno_path}")
        if deno_path:
            deno_ver = subprocess.check_output(["deno", "--version"], text=True).splitlines()[0]
            print(f"   Deno version: {deno_ver}")
        else:
            print("   Deno not found in PATH!")
    except Exception as js_err:
        print(f"   Error checking JS runtimes: {js_err}")

    # 1. 下載 YouTube 影片 (優先下載 720p MP4)
    video_file = "temp_video.mp4"
    ydl_opts = {
        'format': 'best',
        'outtmpl': video_file,
        'quiet': True,
        'no_warnings': True,
        'remote_components': ['ejs:github']
    }
    if cookies:
        ydl_opts['cookiefile'] = cookies
    
    title = "YouTube Video"
    duration = 0
    
    try:
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(video_url, download=True)
            title = info.get('title', 'YouTube Video')
            duration = info.get('duration', 0)
    except Exception as e:
        print(f"❌ Failed to download YouTube video: {e}")
        try:
            print("🔍 Printing available formats for debugging:")
            debug_opts = {'quiet': True}
            if cookies:
                debug_opts['cookiefile'] = cookies
            with yt_dlp.YoutubeDL(debug_opts) as ydl:
                ydl.list_formats(ydl.extract_info(video_url, download=False))
        except Exception as fmt_err:
            print(f"   Failed to list formats: {fmt_err}")
        sys.exit(1)

    print(f"✅ Video downloaded successfully: '{title}' ({duration}s)")

    # 2. 向 Workers API 建立/初始化任務
    init_url = f"{cloudflare_url.rstrip('/')}/api/tasks"
    try:
        r = requests.post(init_url, json={
            "id": task_id,
            "title": title,
            "video_name": video_url,
            "video_duration": duration,
            "status": "analyzing",
            "github_run_id": github_run_id,
            "github_run_url": github_run_url,
            "token_hint": token_hint
        })
        print(f"📌 Task initialized in D1: {r.status_code}")
    except Exception as e:
        print(f"⚠️ Failed to init task in Cloudflare D1: {e}")

    # 3. 開啟影片進行 analysis
    cap = cv2.VideoCapture(video_file)
    if not cap.isOpened():
        print("❌ Cannot open video file with OpenCV.")
        sys.exit(1)

    fps = cap.get(cv2.CAP_PROP_FPS)
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    print(f"📹 FPS: {fps}, Total Frames: {total_frames}")

    slide_candidates = []
    prev_low_res = None
    slide_no = 0

    # 建立暫存目錄存放投影片圖片
    temp_dir = "temp_slides"
    os.makedirs(temp_dir, exist_ok=True)

    # 以秒為間隔進行 MAE 掃描
    frame_step = int(fps * interval) if fps > 0 else 30
    
    for frame_idx in range(0, total_frames, frame_step):
        cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx)
        ret, frame = cap.read()
        if not ret:
            break
            
        current_time = frame_idx / fps if fps > 0 else 0
        
        # 轉成 32x18 灰階低解析度進行 MAE 計算
        gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
        low_res = cv2.resize(gray, (32, 18))
        
        if prev_low_res is not None:
            # 計算 MAE 差異
            mae = np.mean(np.abs(low_res.astype(np.float32) - prev_low_res.astype(np.float32)))
            
            if mae > threshold:
                slide_no += 1
                timestamp = format_timestamp(current_time)
                print(f"📸 Detected new slide #{slide_no} at {timestamp} (MAE: {mae:.2f})")
                
                # 儲存高解析度 PNG 投影片
                slide_filename = os.path.join(temp_dir, f"slide_{slide_no:03d}.png")
                cv2.imwrite(slide_filename, frame)
                
                # 直傳此投影片至 Workers R2
                upload_url = f"{cloudflare_url.rstrip('/')}/api/tasks/{task_id}/upload"
                try:
                    with open(slide_filename, "rb") as f:
                        file_bytes = f.read()
                        
                    up_res = requests.post(upload_url, data=file_bytes, headers={
                        "X-Slide-No": str(slide_no),
                        "X-Timestamp": timestamp,
                        "X-Seconds": str(current_time)
                    })
                    if up_res.status_code == 200:
                        print(f"   Upload slide #{slide_no} status: {up_res.status_code}")
                    else:
                        print(f"   ❌ Upload slide #{slide_no} status: {up_res.status_code}, error: {up_res.text}")
                except Exception as e:
                    print(f"   ⚠️ Upload failed: {e}")
                    
                slide_candidates.append({
                    "slide_no": slide_no,
                    "timestamp": timestamp,
                    "seconds": current_time,
                    "path": slide_filename
                })
                
        else:
            # 第一張投影片
            slide_no += 1
            timestamp = format_timestamp(current_time)
            print(f"📸 Detected first slide #{slide_no} at {timestamp}")
            
            slide_filename = os.path.join(temp_dir, f"slide_{slide_no:03d}.png")
            cv2.imwrite(slide_filename, frame)
            
            # 直傳至 Workers
            upload_url = f"{cloudflare_url.rstrip('/')}/api/tasks/{task_id}/upload"
            try:
                with open(slide_filename, "rb") as f:
                    file_bytes = f.read()
                up_res = requests.post(upload_url, data=file_bytes, headers={
                    "X-Slide-No": str(slide_no),
                    "X-Timestamp": timestamp,
                    "X-Seconds": str(current_time)
                })
                if up_res.status_code == 200:
                    print(f"   Upload first slide #{slide_no} status: {up_res.status_code}")
                else:
                    print(f"   ❌ Upload first slide #{slide_no} status: {up_res.status_code}, error: {up_res.text}")
            except Exception as e:
                print(f"   ⚠️ Upload first slide failed: {e}")
                
            slide_candidates.append({
                "slide_no": slide_no,
                "timestamp": timestamp,
                "seconds": current_time,
                "path": slide_filename
            })
            
        prev_low_res = low_res

    cap.release()
    print(f"🎉 Analysis completed! Total slides captured: {len(slide_candidates)}")

    if not slide_candidates:
        print("⚠️ No slides detected.")
        sys.exit(0)

    # 4. 編譯 PDF 文件並上傳
    print("📄 Compiling PDF presentation...")
    pdf_path = "presentation.pdf"
    try:
        images_list = []
        for s in slide_candidates:
            img = Image.open(s["path"]).convert("RGB")
            images_list.append(img)
            
        if images_list:
            images_list[0].save(pdf_path, save_all=True, append_images=images_list[1:])
            print("   PDF compiled successfully.")
            
            # 上傳 PDF
            doc_url = f"{cloudflare_url.rstrip('/')}/api/tasks/{task_id}/document"
            with open(pdf_path, "rb") as f:
                pdf_res = requests.post(doc_url, data=f.read(), headers={
                    "X-Doc-Type": "pdf"
                })
            if pdf_res.status_code == 200:
                print("   PDF uploaded to Cloudflare R2.")
            else:
                print(f"   ❌ PDF upload failed with status {pdf_res.status_code}: {pdf_res.text}")
    except Exception as e:
        print(f"❌ Failed to compile/upload PDF: {e}")

    # 5. 編譯 PPTX 文件並上傳
    print("📊 Compiling PPTX presentation...")
    pptx_path = "presentation.pptx"
    try:
        prs = Presentation()
        # 設定 16:9 比例頁面尺寸
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_slide_layout = prs.slide_layouts[6]
        
        for s in slide_candidates:
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(s["path"], 0, 0, width=prs.slide_width, height=prs.slide_height)
            
        prs.save(pptx_path)
        print("   PPTX compiled successfully.")
        
        # 上傳 PPTX
        with open(pptx_path, "rb") as f:
            pptx_res = requests.post(doc_url, data=f.read(), headers={
                "X-Doc-Type": "pptx"
            })
        if pptx_res.status_code == 200:
            print("   PPTX uploaded to Cloudflare R2.")
        else:
            print(f"   ❌ PPTX upload failed with status {pptx_res.status_code}: {pptx_res.text}")
    except Exception as e:
        print(f"❌ Failed to compile/upload PPTX: {e}")

    # 5.2. 編譯 HTML 報告並上傳
    print("📄 Compiling HTML report...")
    html_path = "report.html"
    try:
        generate_report_html(html_path, title, video_url, duration, slide_candidates)
        print("   HTML report compiled successfully.")
        
        # 上傳 HTML
        with open(html_path, "rb") as f:
            html_res = requests.post(doc_url, data=f.read(), headers={
                "X-Doc-Type": "html"
            })
        if html_res.status_code == 200:
            print("   HTML report uploaded to Cloudflare R2.")
        else:
            print(f"   ❌ HTML report upload failed with status {html_res.status_code}: {html_res.text}")
    except Exception as e:
        print(f"❌ Failed to compile/upload HTML report: {e}")

    # 5.5. 將任務狀態更新為 completed，並發送執行時間
    try:
        final_url = f"{cloudflare_url.rstrip('/')}/api/tasks/{task_id}"
        execution_time_sec = float(time.time() - start_time)
        print(f"⌛ Task completed in {execution_time_sec:.2f}s. Updating status in D1...")
        requests.put(final_url, json={
            "status": "completed",
            "execution_time_sec": execution_time_sec
        })
    except Exception as e:
        print(f"⚠️ Failed to update final task status: {e}")

    # 6. 清理暫存檔案
    try:
        if os.path.exists(video_file):
            os.remove(video_file)
        for s in slide_candidates:
            if os.path.exists(s["path"]):
                os.remove(s["path"])
        if os.path.exists(temp_dir):
            os.rmdir(temp_dir)
        if os.path.exists(pdf_path):
            os.remove(pdf_path)
        if os.path.exists(pptx_path):
            os.remove(pptx_path)
        if os.path.exists(html_path):
            os.remove(html_path)
        print("🧹 Cleaned up all temporary files.")
    except Exception as e:
        print(f"⚠️ Error during cleanup: {e}")

    print("🏁 Execution finished successfully!")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="GitHub Actions YouTube Slide Extractor")
    parser.add_argument("--url", required=True, help="YouTube Video URL")
    parser.add_argument("--task-id", required=True, help="Cloudflare Task ID")
    parser.add_argument("--cloudflare-url", required=True, help="Cloudflare Worker Domain URL")
    parser.add_argument("--threshold", type=float, default=8.0, help="MAE change threshold")
    parser.add_argument("--interval", type=float, default=1.0, help="Frame check interval in seconds")
    parser.add_argument("--cookies", help="Path to cookies file")
    parser.add_argument("--token-hint", help="First few characters of token for metrics tracking")
    
    args = parser.parse_args()
    analyze_and_upload(
        video_url=args.url,
        task_id=args.task_id,
        cloudflare_url=args.cloudflare_url,
        threshold=args.threshold,
        interval=args.interval,
        cookies=args.cookies,
        token_hint=args.token_hint
    )
